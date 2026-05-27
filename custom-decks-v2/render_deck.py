"""Render a custom prospect deck for a prospect, upload to Drive, return links."""
import os
import sys, os, subprocess, datetime, argparse, json, time, random, re
from gdrive_peace import get_creds
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Per-token character budgets — enforced before render to prevent layout overflow.
TOKEN_BUDGETS = {
    "COMPANY_NAME":       40,
    "FIRST_NAME":         20,
    "SHIFT_HEADLINE":     42,
    "SHIFT_PUNCHLINE":    42,
    "SHIFT_BODY":         260,   # may contain <br>; we measure visible chars
    "SHIFT_STAT1_NUM":    8,
    "SHIFT_STAT1_HEAD":   25,
    "SHIFT_STAT1_BODY":   100,
    "SHIFT_STAT2_NUM":    8,
    "SHIFT_STAT2_HEAD":   25,
    "SHIFT_STAT2_BODY":   100,
    "SHIFT_STAT3_NUM":    8,
    "SHIFT_STAT3_HEAD":   25,
    "SHIFT_STAT3_BODY":   100,
    "SHIFT_CLOSE":        180,
    "COMPETE_HEAD1":      42,
    "COMPETE_HEAD2":      42,
    "COMPETE_INTRO":      240,
    "COMPETITOR_NAME":    30,
    "COMPETITOR_TAGLINE": 60,
    "COMPETITOR_LINE1":   65,
    "COMPETITOR_LINE2":   65,
    "COMPETITOR_LINE3":   65,
    "COMPANY_TAGLINE":    60,
    "COMPANY_LINE1":      65,
    "COMPANY_LINE2":      65,
    "COMPANY_LINE3":      65,
    "WHY_OUTBOUND_FITS":  200,
    "ICP_LABEL":          50,
    "THEIR_TARGET_TITLE": 40,
    "ASK_HEAD1":          40,
    "ASK_HEAD2":          50,
    "ASK_INTRO":          240,
}

def _visible_len(s):
    """Approximate visible char count (strip simple HTML tags)."""
    return len(re.sub(r'<[^>]+>', '', s))

def trim_to_budget(text, budget):
    """If text exceeds budget (visible chars), truncate at a sentence/word boundary
    and append an ellipsis. Preserves HTML tags reasonably."""
    if _visible_len(text) <= budget:
        return text
    plain = re.sub(r'<[^>]+>', ' ', text)
    cut = plain[:budget]
    # Try to break on sentence boundary first, then space.
    for sep in ['. ', '! ', '? ', '; ', ' — ', ', ', ' ']:
        i = cut.rfind(sep)
        if i >= int(budget * 0.6):
            cut = cut[:i + len(sep)].rstrip()
            break
    return cut.rstrip(' .,;:—') + '…'

HEADER_PAIRS = [
    ('SHIFT_HEADLINE',  'SHIFT_PUNCHLINE'),
    ('COMPETE_HEAD1',   'COMPETE_HEAD2'),
    ('ASK_HEAD1',       'ASK_HEAD2'),
]

def enforce_header_style(tokens):
    """Strip ALL trailing punctuation from the first line of every compound header.
    Per Nelson: no period, no dash, no comma — just words. The second line is
    bolded by the template and renders in brand color."""
    for top, bot in HEADER_PAIRS:
        v = tokens.get(top, '')
        if not isinstance(v, str): continue
        s = v.strip()
        if not s: continue
        # Strip ANY trailing punctuation + whitespace.
        s = re.sub(r'[\s.,;:!\?\-–—]+$', '', s)
        if s:
            tokens[top] = s
        b = (tokens.get(bot, '') or '').strip()
        if b:
            tokens[bot] = b

def enforce_budgets(tokens, verbose=False):
    """Trim any token that exceeds its char budget. Returns list of trimmed keys."""
    trimmed = []
    for key, budget in TOKEN_BUDGETS.items():
        v = tokens.get(key, '')
        if not isinstance(v, str): continue
        if _visible_len(v) > budget:
            tokens[key] = trim_to_budget(v, budget)
            trimmed.append(key)
            if verbose:
                print(f"  trimmed {key}: {_visible_len(v)}→{_visible_len(tokens[key])} chars")
    return trimmed

# Tokens that render as multi-line paragraphs. The rule (Lara Acosta style):
# 1. Every sentence on its OWN line (split on . ! ?).
# 2. If a single sentence is too long for the per-line budget, word-wrap it.
LINE_WRAP_TOKENS = {
    'SHIFT_BODY':        72,
    'COMPETE_INTRO':     76,
    'WHY_OUTBOUND_FITS': 76,
    'SHIFT_CLOSE':       80,
    'ASK_INTRO':         76,
    'SHIFT_STAT1_BODY':  38,
    'SHIFT_STAT2_BODY':  38,
    'SHIFT_STAT3_BODY':  38,
}

_SENT_RX = re.compile(r'(?<=[.!?])\s+(?=[A-Z0-9])')

def _word_wrap(sentence, per_line):
    """Wrap a single long sentence at word boundaries; returns <br>-joined chunks."""
    if _visible_len(sentence) <= per_line:
        return sentence
    out = []
    rest = sentence.strip()
    while _visible_len(rest) > per_line:
        head = rest[:per_line + 1]
        idx = head.rfind(' ')
        if idx <= int(per_line * 0.4):
            idx = per_line  # hard break only if no decent word boundary
        out.append(rest[:idx].rstrip())
        rest = rest[idx:].lstrip()
    out.append(rest)
    return '<br>'.join(out)

def enforce_line_wrap(tokens, verbose=False):
    """For paragraph-style tokens, split on sentence boundaries first (every
    sentence on its own line), then word-wrap any sentence that exceeds the
    per-line budget."""
    wrapped = []
    for key, per_line in LINE_WRAP_TOKENS.items():
        v = tokens.get(key, '')
        if not isinstance(v, str) or not v.strip(): continue
        original = v
        # Existing <br> already mark sentence-ish boundaries; respect them.
        chunks = re.split(r'<br\s*/?>', v)
        out_chunks = []
        for c in chunks:
            c = c.strip()
            if not c: continue
            # Split further on real sentence boundaries
            sentences = _SENT_RX.split(c)
            for s in sentences:
                s = s.strip()
                if not s: continue
                out_chunks.append(_word_wrap(s, per_line))
        new = '<br>'.join(out_chunks)
        if new != original:
            tokens[key] = new
            wrapped.append(key)
            if verbose:
                print(f"  wrapped {key}: {original.count('<br>')+1} → {new.count('<br>')+1} lines")
    return wrapped

def with_backoff(fn, *, retries=6, base=5.0, label="drive"):
    """Run a Drive API call with exponential backoff on 429/rateLimit/5xx."""
    for attempt in range(retries):
        try:
            return fn()
        except Exception as e:
            msg = str(e)
            transient = ('429' in msg or 'rateLimit' in msg or
                         'userRateLimit' in msg or '500' in msg or '503' in msg)
            if not transient or attempt == retries - 1:
                raise
            wait = base * (2 ** attempt) + random.uniform(0, 2)
            print(f"  {label} retry {attempt+1}/{retries} after {wait:.1f}s: {msg[:120]}")
            time.sleep(wait)

TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)),"deck_template.md.tpl")
OUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),"rendered")
DRIVE_FOLDER_ID = None  # will be created/discovered

def discover_or_create_folder(drive, name=os.environ.get("DECK_FOLDER_NAME","Prospect Decks")):
    q = f"name='{name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
    res = with_backoff(
        lambda: drive.files().list(q=q, fields="files(id,name)").execute(),
        label="folder.list"
    ).get("files", [])
    if res:
        return res[0]["id"]
    folder = with_backoff(lambda: drive.files().create(
        body={"name": name, "mimeType": "application/vnd.google-apps.folder"},
        fields="id"
    ).execute(), label="folder.create")
    return folder["id"]

def render(tokens):
    with open(TEMPLATE) as f:
        md = f.read()
    for k, v in tokens.items():
        md = md.replace(f"{{{{{k}}}}}", v)
    slug = tokens["COMPANY_NAME"].lower().replace(" ", "_").replace(".", "")
    md_path = f"{OUT_DIR}/{slug}.md"
    pdf_path = f"{OUT_DIR}/{slug}.pdf"
    pptx_path = f"{OUT_DIR}/{slug}.pptx"
    with open(md_path, "w") as f:
        f.write(md)
    subprocess.run(
        ["marp", md_path, "--pdf", "--allow-local-files", "-o", pdf_path],
        check=True, capture_output=True
    )
    subprocess.run(
        ["marp", md_path, "--pptx", "--allow-local-files", "-o", pptx_path],
        check=True, capture_output=True
    )
    return md_path, pdf_path, pptx_path, slug

def find_cta_bbox(pdf_path):
    """Render the last PDF page at 96dpi and return the CTA orange-button bbox
    in pixel coords (x1, y1, x2, y2). Marp PDF is 1280x720 at 96dpi."""
    out_prefix = pdf_path.replace(".pdf", "_lastpage")
    # determine page count
    info = subprocess.run(["pdfinfo", pdf_path], capture_output=True, text=True)
    pages = 1
    for line in info.stdout.splitlines():
        if line.startswith("Pages:"):
            pages = int(line.split()[1])
    subprocess.run(
        ["pdftoppm", "-r", "96", "-f", str(pages), "-l", str(pages),
         pdf_path, out_prefix, "-png"],
        check=True, capture_output=True
    )
    png = f"{out_prefix}-{pages:0{len(str(pages))}d}.png"
    if not os.path.exists(png):
        # pdftoppm zero-pads sometimes
        for cand in [f"{out_prefix}-{pages}.png", f"{out_prefix}-{pages:02d}.png",
                     f"{out_prefix}-{pages:03d}.png"]:
            if os.path.exists(cand):
                png = cand; break
    im = Image.open(png).convert("RGB")
    W, H = im.size
    px = im.load()
    def near(c, t, tol=50):
        return all(abs(c[i]-t[i]) < tol for i in range(3))
    rows = {}
    # CTA sits in the lower portion of the slide
    for y in range(int(H*0.65), int(H*0.92)):
        cnt = 0; minx = W; maxx = 0
        for x in range(W):
            c = px[x, y]
            if (near(c, (253, 82, 53)) or near(c, (254, 129, 30)) or
                    near(c, (253, 44, 72))):
                cnt += 1
                if x < minx: minx = x
                if x > maxx: maxx = x
        if cnt > 30:
            rows[y] = (cnt, minx, maxx)
    if not rows:
        return None
    ys = sorted(rows.keys())
    x1 = min(r[1] for r in rows.values())
    x2 = max(r[2] for r in rows.values())
    return (x1, ys[0], x2, ys[-1], W, H)

def patch_cta_hotspot(pptx_path, pdf_path, booking_url):
    """Drop a transparent clickable rectangle over the CTA button on the
    last slide, since Marp PPTX flattens everything to a background image."""
    bbox = find_cta_bbox(pdf_path)
    if not bbox:
        print("warn: CTA bbox not found; skipping hotspot")
        return
    x1, y1, x2, y2, img_w, img_h = bbox
    p = Presentation(pptx_path)
    slide_w, slide_h = p.slide_width, p.slide_height
    sx = slide_w / img_w
    sy = slide_h / img_h
    pad = 4  # px pad so the rectangle fully covers the visible button
    left = int((x1 - pad) * sx)
    top = int((y1 - pad) * sy)
    width = int((x2 - x1 + 2*pad) * sx)
    height = int((y2 - y1 + 2*pad) * sy)
    last = p.slides[-1]
    shp = last.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.background()
    shp.line.fill.background()
    shp.click_action.hyperlink.address = booking_url
    shp.name = "CTA_HOTSPOT"
    p.save(pptx_path)
    print(f"hotspot: x={x1}-{x2} y={y1}-{y2} -> {booking_url}")

def upload(drive, folder_id, pdf_path, pptx_path, company_name,
           slides_id=None, pdf_id=None):
    pdf_media = MediaFileUpload(pdf_path, mimetype="application/pdf", resumable=False)
    pptx_media = MediaFileUpload(
        pptx_path,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False,
    )
    if slides_id:
        slides_res = with_backoff(lambda: drive.files().update(
            fileId=slides_id, media_body=pptx_media,
            fields="id,webViewLink"
        ).execute(), label="slides.update")
    else:
        slides_res = with_backoff(lambda: drive.files().create(
            body={"name": f"{company_name} deck",
                  "mimeType": "application/vnd.google-apps.presentation",
                  "parents": [folder_id]},
            media_body=pptx_media, fields="id,webViewLink"
        ).execute(), label="slides.create")
    if pdf_id:
        pdf_res = with_backoff(lambda: drive.files().update(
            fileId=pdf_id, media_body=pdf_media,
            fields="id,webViewLink"
        ).execute(), label="pdf.update")
    else:
        pdf_res = with_backoff(lambda: drive.files().create(
            body={"name": f"{company_name} deck.pdf",
                  "parents": [folder_id]},
            media_body=pdf_media, fields="id,webViewLink"
        ).execute(), label="pdf.create")
    for fid in [pdf_res["id"], slides_res["id"]]:
        try:
            with_backoff(lambda fid=fid: drive.permissions().create(
                fileId=fid, body={"role": "reader", "type": "anyone"},
                fields="id"
            ).execute(), label=f"perm.{fid[:8]}")
        except Exception:
            pass  # already shared
    # Publish-to-web on the Slides file so the CTA button is clickable in present mode
    try:
        revs = with_backoff(lambda: drive.revisions().list(
            fileId=slides_res["id"], fields="revisions(id)"
        ).execute(), label="revs.list").get("revisions", [])
        if revs:
            with_backoff(lambda: drive.revisions().update(
                fileId=slides_res["id"], revisionId=revs[-1]["id"],
                body={"published": True, "publishedOutsideDomain": True,
                      "publishAuto": True}
            ).execute(), label="revs.publish")
    except Exception as e:
        print(f"publish-to-web warning: {e}")
    # NOTE: /pub URL is unreliable on this Workspace; /preview works publicly when file
    # has "anyone with link" permission. Caller (cron) should prefer /preview for sharing.
    pub_link = (f"https://docs.google.com/presentation/d/{slides_res['id']}"
                "/preview")
    return pdf_res["webViewLink"], slides_res["webViewLink"], pub_link

def main():
    p = argparse.ArgumentParser()
    p.add_argument("--tokens", required=True, help="JSON file with token values")
    p.add_argument("--slides-id", help="Existing Google Slides file ID to update in-place")
    p.add_argument("--pdf-id", help="Existing Drive PDF file ID to update in-place")
    args = p.parse_args()
    with open(args.tokens) as f:
        tokens = json.load(f)
    tokens.setdefault("DATE", datetime.date.today().strftime("%B %Y"))
    tokens.setdefault("LOGO_HEIGHT", "40px")
    # Title mark: always the centered RC mark. Per Nelson, prospect logos are NOT
    # used on the title slide, even when available.
    tokens["TITLE_MARK_BLOCK"] = (
        '<div class="title-row" style="justify-content: center;">'
        '<div class="title-mark" style="width: 140px; height: 140px; margin-bottom: 0;"></div>'
        '</div>'
    )
    tokens.setdefault("ASK_HEAD1", "30 minutes.")
    tokens.setdefault("ASK_HEAD2", "That's the whole ask.")
    tokens.setdefault("ASK_INTRO",
        f"Worst case: a coverage benchmark vs other {tokens.get('ICP_LABEL','')}.<br>"
        f"Best case: {tokens.get('COMPANY_NAME','')} owns the supplier graph before "
        "the category gets named.")
    # Enforce header style first so the trailing " —" is counted in budgets
    enforce_header_style(tokens)
    # Enforce per-token char budgets so nothing overflows the slide
    trimmed = enforce_budgets(tokens, verbose=True)
    if trimmed:
        print(f"trimmed tokens: {', '.join(trimmed)}")
    # Enforce per-line wrap so no sentence breaks mid-line on the slide
    wrapped = enforce_line_wrap(tokens, verbose=True)
    if wrapped:
        print(f"wrapped tokens: {', '.join(wrapped)}")
    # Always persist post-processed tokens — header style, trimming, and wrapping
    # may all mutate things and downstream self-checks should see the same content.
    try:
        with open(args.tokens, "w") as f:
            json.dump(tokens, f, indent=2)
    except Exception:
        pass
    md_path, pdf_path, pptx_path, slug = render(tokens)
    print(f"rendered: {md_path}")
    print(f"pdf: {pdf_path}")
    print(f"pptx: {pptx_path}")
    if tokens.get("BOOKING_LINK"):
        patch_cta_hotspot(pptx_path, pdf_path, tokens["BOOKING_LINK"])
    drive = build("drive", "v3", credentials=get_creds())
    folder_id = discover_or_create_folder(drive)
    print(f"folder: {folder_id}")
    pdf_link, slides_link, pub_link = upload(drive, folder_id, pdf_path,
                                              pptx_path, tokens["COMPANY_NAME"],
                                              slides_id=args.slides_id,
                                              pdf_id=args.pdf_id)
    print(f"PDF_LINK: {pdf_link}")
    print(f"SLIDES_LINK: {slides_link}")
    print(f"PUB_LINK: {pub_link}")

if __name__ == "__main__":
    main()
